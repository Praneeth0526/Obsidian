import sys, os
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))
"""
test_file_generator.py — Programmatic test file generation for E2E pipeline QA.

Generates files in 6 formats × 3 size variants + edge cases.
Each file contains realistic, varied content designed to exercise
the full Tika extraction → LangChain chunking pipeline.

Test-only dependencies:
    pip install reportlab python-docx python-pptx Pillow

Usage:
    from test_file_generator import generate_all, TestFile
    files = generate_all()
    for f in files:
        print(f"{f.test_case_id}: {f.file_format} {f.size_variant} "
              f"({f.file_size_bytes:,} bytes)")
"""

import io
import logging
import os
import random
from dataclasses import dataclass
from typing import List, Optional

from PIL import Image, ImageDraw, ImageFilter

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Optional dependencies — graceful degradation
# ---------------------------------------------------------------------------

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    from docx import Document as DocxDocument
    from docx.shared import Pt as DocxPt

    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    from pptx.util import Pt as PptxPt

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# ---------------------------------------------------------------------------
# TestFile dataclass
# ---------------------------------------------------------------------------

@dataclass
class TestFile:
    """A generated test file with metadata."""

    name: str                # e.g., "pdf_small.pdf"
    test_case_id: str        # e.g., "pdf_small"
    file_format: str         # e.g., "PDF"
    size_variant: str        # "small", "medium", "large", or "edge"
    content_type: str        # MIME type
    file_bytes: bytes        # raw file content
    description: str         # human-readable description
    is_edge_case: bool = False
    expected_failure: bool = False   # True for corrupt/empty files

    @property
    def file_size_bytes(self) -> int:
        return len(self.file_bytes)


# ---------------------------------------------------------------------------
# Sample text content — realistic and varied
# ---------------------------------------------------------------------------

SAMPLE_PARAGRAPHS = [
    (
        "The Apache Tika toolkit detects and extracts metadata and text "
        "from over a thousand different file types. All of these file types "
        "can be parsed through a single interface, making Tika useful for "
        "search engine indexing, content analysis, translation, and much more. "
        "Tika exposes a REST API that accepts file bytes and returns extracted "
        "plain text along with structured metadata such as author, creation "
        "date, page count, and content type."
    ),
    (
        "Natural language processing is a subfield of linguistics, computer "
        "science, and artificial intelligence concerned with the interactions "
        "between computers and human language. It focuses on how to program "
        "computers to process and analyze large amounts of natural language "
        "data. Applications include machine translation, sentiment analysis, "
        "named entity recognition, and text summarization."
    ),
    (
        "Vector databases store data as high-dimensional vectors, which are "
        "mathematical representations of features or attributes. Each vector "
        "has a certain number of dimensions, which can range from tens to "
        "thousands, depending on the complexity and granularity of the data. "
        "Cosine similarity and dot product are common distance metrics used "
        "to find nearest neighbors in vector space."
    ),
    (
        "Apache Kafka is a distributed event streaming platform used by "
        "thousands of companies for high-performance data pipelines, streaming "
        "analytics, data integration, and mission-critical applications. "
        "It combines three key capabilities: publishing and subscribing to "
        "streams of events, storing streams of events durably and reliably, "
        "and processing streams of events as they occur."
    ),
    (
        "MinIO is a high-performance, S3-compatible object store built for "
        "large-scale AI and ML workloads, data lakes, and database workloads. "
        "It is software-defined and runs on any cloud or on-premises "
        "infrastructure. MinIO supports bucket notifications that can trigger "
        "events on object creation, deletion, and access, making it ideal "
        "for event-driven architectures."
    ),
    (
        "OpenSearch is a community-driven, open-source search and analytics "
        "suite derived from Apache Lucene. It supports full-text search, "
        "structured search, analytics, and k-nearest neighbor vector search. "
        "The k-NN plugin enables approximate nearest neighbor search using "
        "HNSW graphs, supporting use cases like semantic search, "
        "recommendation engines, and anomaly detection."
    ),
    (
        "Sentence transformers are neural network models that map sentences "
        "and paragraphs to dense vector representations in a shared embedding "
        "space. The all-MiniLM-L6-v2 model produces 384-dimensional vectors "
        "and is optimized for semantic similarity tasks. It processes text "
        "through a BERT-like architecture with 6 transformer layers and "
        "achieves a good balance between quality and inference speed."
    ),
    (
        "The ingestion pipeline processes uploaded files through multiple "
        "stages: text extraction via Apache Tika, recursive character-based "
        "chunking via LangChain, embedding generation via SentenceTransformers, "
        "and storage in OpenSearch for hybrid search. Each chunk carries "
        "positional metadata including character offsets and chunk indices "
        "to enable passage highlighting in search results."
    ),
    (
        "Retrieval-augmented generation combines the power of large language "
        "models with external knowledge retrieval. By fetching relevant "
        "documents from a vector database based on semantic similarity, "
        "the system provides more accurate, grounded, and verifiable "
        "responses. This approach reduces hallucination and allows the "
        "model to access up-to-date information beyond its training data."
    ),
    (
        "Microservices architecture structures an application as a collection "
        "of loosely coupled services that implement specific business "
        "capabilities. Each service can be developed, deployed, and scaled "
        "independently. Communication between services typically occurs "
        "through REST APIs, message queues like Kafka, or gRPC. This "
        "architecture enables teams to work autonomously and adopt "
        "different technology stacks for different services."
    ),
    (
        "Docker containers package applications with their dependencies "
        "into standardized units for development, shipment, and deployment. "
        "Docker Compose defines multi-container applications using a YAML "
        "file, specifying services, networks, volumes, and dependencies. "
        "Health checks ensure containers are ready before dependent services "
        "start, preventing race conditions during orchestration."
    ),
    (
        "Text chunking is a critical preprocessing step in RAG pipelines. "
        "The recursive character text splitter tries to split text by "
        "progressively smaller separators: first by double newlines "
        "(paragraphs), then single newlines, then spaces, and finally "
        "individual characters. Overlapping chunks ensure that context "
        "at chunk boundaries is preserved for downstream embedding and "
        "retrieval, preventing information loss at split points."
    ),
]

SLIDE_TOPICS = [
    ("Introduction to Data Pipelines", "Data pipelines automate the flow of data from source systems to destinations. They handle extraction, transformation, and loading (ETL) of data at scale."),
    ("Architecture Overview", "The system uses a microservices architecture with Kafka for event streaming, MinIO for object storage, and OpenSearch for search and analytics."),
    ("Text Extraction with Tika", "Apache Tika extracts text and metadata from uploaded documents. It supports PDF, DOCX, PPTX, TXT, and many other formats through a unified REST API."),
    ("Chunking Strategy", "Documents are split into overlapping chunks of 512 characters with 50-character overlap. This preserves context at boundaries for better retrieval."),
    ("Embedding Models", "The all-MiniLM-L6-v2 model generates 384-dimensional vectors for text. Nomic-embed-vision-v1.5 handles image embeddings in the same vector space."),
    ("Vector Search", "OpenSearch k-NN plugin enables approximate nearest neighbor search using HNSW graphs. Cosine similarity measures semantic closeness between vectors."),
    ("Hybrid Search", "Combining BM25 keyword search with k-NN vector search improves retrieval quality. Results are fused using reciprocal rank fusion."),
    ("Image Processing", "Images are preprocessed to 224x224 RGB JPEG format before embedding. The ImageHandler normalizes color modes and resize dimensions."),
    ("Error Handling", "The pipeline handles corrupt files, empty documents, and network failures gracefully. Failed messages are routed to a dead-letter queue."),
    ("Performance Metrics", "Key metrics include extraction time, chunking throughput, embedding latency, and end-to-end pipeline latency per document."),
    ("Scaling Considerations", "Kafka partitions enable parallel consumption. The model server can be horizontally scaled behind a load balancer."),
    ("Security and Access Control", "MinIO supports IAM policies for bucket-level access control. OpenSearch provides role-based access to indices and documents."),
    ("Monitoring and Alerting", "Structured logging with structlog enables centralized log aggregation. Health check endpoints provide liveness and readiness probes."),
    ("Deployment Strategy", "Services are deployed as Docker containers orchestrated by Docker Compose locally and Kubernetes in production environments."),
    ("Future Roadmap", "Planned features include OCR support for scanned documents, multi-language embedding models, and real-time streaming analytics."),
]


def _repeat_paragraphs(count: int) -> list[str]:
    """Repeat and number sample paragraphs to produce *count* unique paragraphs."""
    result = []
    for i in range(count):
        base = SAMPLE_PARAGRAPHS[i % len(SAMPLE_PARAGRAPHS)]
        result.append(f"[Section {i + 1}] {base}")
    return result


def _repeat_slides(count: int) -> list[tuple[str, str]]:
    """Repeat and number slide topics to produce *count* unique slides."""
    result = []
    for i in range(count):
        title, content = SLIDE_TOPICS[i % len(SLIDE_TOPICS)]
        result.append((f"Slide {i + 1}: {title}", content))
    return result


# ---------------------------------------------------------------------------
# PDF Generation (requires reportlab)
# ---------------------------------------------------------------------------

def _build_pdf(paragraphs: list[str], pagesize=None) -> bytes:
    """Build a PDF from a list of paragraph strings using reportlab."""
    if not HAS_REPORTLAB:
        raise ImportError(
            "reportlab is required for PDF generation. "
            "Install with: pip install reportlab"
        )

    if pagesize is None:
        pagesize = letter

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=pagesize)
    styles = getSampleStyleSheet()
    normal = styles["Normal"]

    story = []
    for p in paragraphs:
        story.append(Paragraph(p, normal))
        story.append(Spacer(1, 12))

    doc.build(story)
    return buffer.getvalue()


def generate_pdf_small() -> TestFile:
    """1 page, 3 sentences."""
    paragraphs = SAMPLE_PARAGRAPHS[:3]
    return TestFile(
        name="pdf_small.pdf",
        test_case_id="pdf_small",
        file_format="PDF",
        size_variant="small",
        content_type="application/pdf",
        file_bytes=_build_pdf(paragraphs),
        description="1-page PDF with 3 paragraphs (~900 chars of content)",
    )


def generate_pdf_medium() -> TestFile:
    """~10 pages of varied paragraphs."""
    paragraphs = _repeat_paragraphs(50)
    return TestFile(
        name="pdf_medium.pdf",
        test_case_id="pdf_medium",
        file_format="PDF",
        size_variant="medium",
        content_type="application/pdf",
        file_bytes=_build_pdf(paragraphs),
        description="~10-page PDF with 50 paragraphs (~15,000 chars)",
    )


def generate_pdf_large() -> TestFile:
    """~60 pages of dense text."""
    paragraphs = _repeat_paragraphs(300)
    return TestFile(
        name="pdf_large.pdf",
        test_case_id="pdf_large",
        file_format="PDF",
        size_variant="large",
        content_type="application/pdf",
        file_bytes=_build_pdf(paragraphs),
        description="~60-page PDF with 300 paragraphs (~90,000 chars)",
    )


# ---------------------------------------------------------------------------
# DOCX Generation (requires python-docx)
# ---------------------------------------------------------------------------

def _build_docx(paragraphs: list[str], with_headings: bool = False) -> bytes:
    """Build a DOCX from a list of paragraph strings."""
    if not HAS_DOCX:
        raise ImportError(
            "python-docx is required for DOCX generation. "
            "Install with: pip install python-docx"
        )

    doc = DocxDocument()

    for i, p in enumerate(paragraphs):
        if with_headings and i % 5 == 0:
            doc.add_heading(f"Chapter {i // 5 + 1}", level=2)
        doc.add_paragraph(p)

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def generate_docx_small() -> TestFile:
    """1 paragraph."""
    paragraphs = SAMPLE_PARAGRAPHS[:2]
    return TestFile(
        name="docx_small.docx",
        test_case_id="docx_small",
        file_format="DOCX",
        size_variant="small",
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        file_bytes=_build_docx(paragraphs),
        description="Small DOCX with 2 paragraphs (~600 chars)",
    )


def generate_docx_medium() -> TestFile:
    """~15 pages with headings."""
    paragraphs = _repeat_paragraphs(60)
    return TestFile(
        name="docx_medium.docx",
        test_case_id="docx_medium",
        file_format="DOCX",
        size_variant="medium",
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        file_bytes=_build_docx(paragraphs, with_headings=True),
        description="~15-page DOCX with 60 paragraphs and headings (~18,000 chars)",
    )


def generate_docx_large() -> TestFile:
    """~50 pages with tables and headings."""
    paragraphs = _repeat_paragraphs(250)
    return TestFile(
        name="docx_large.docx",
        test_case_id="docx_large",
        file_format="DOCX",
        size_variant="large",
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        file_bytes=_build_docx(paragraphs, with_headings=True),
        description="~50-page DOCX with 250 paragraphs and headings (~75,000 chars)",
    )


# ---------------------------------------------------------------------------
# TXT Generation
# ---------------------------------------------------------------------------

def generate_txt_small() -> TestFile:
    """~200 chars of plain text."""
    text = " ".join(SAMPLE_PARAGRAPHS[:1])
    return TestFile(
        name="txt_small.txt",
        test_case_id="txt_small",
        file_format="TXT",
        size_variant="small",
        content_type="text/plain",
        file_bytes=text.encode("utf-8"),
        description=f"Small TXT with ~{len(text)} chars",
    )


def generate_txt_medium() -> TestFile:
    """~50 KB of text."""
    paragraphs = _repeat_paragraphs(150)
    text = "\n\n".join(paragraphs)
    return TestFile(
        name="txt_medium.txt",
        test_case_id="txt_medium",
        file_format="TXT",
        size_variant="medium",
        content_type="text/plain",
        file_bytes=text.encode("utf-8"),
        description=f"Medium TXT with ~{len(text):,} chars (~{len(text)//1024} KB)",
    )


def generate_txt_large() -> TestFile:
    """~500 KB of text."""
    paragraphs = _repeat_paragraphs(1500)
    text = "\n\n".join(paragraphs)
    return TestFile(
        name="txt_large.txt",
        test_case_id="txt_large",
        file_format="TXT",
        size_variant="large",
        content_type="text/plain",
        file_bytes=text.encode("utf-8"),
        description=f"Large TXT with ~{len(text):,} chars (~{len(text)//1024} KB)",
    )


# ---------------------------------------------------------------------------
# PNG Generation (uses Pillow — already in requirements.txt)
# ---------------------------------------------------------------------------

def _build_image(
    width: int,
    height: int,
    fmt: str = "PNG",
    text: str = "Test Image",
    add_shapes: bool = False,
    add_gradient: bool = False,
    add_noise: bool = False,
) -> bytes:
    """Build a synthetic image with optional visual complexity."""
    img = Image.new("RGB", (width, height), color=(42, 60, 85))
    draw = ImageDraw.Draw(img)

    if add_gradient:
        for y in range(height):
            r = int(42 + (180 * y / height))
            g = int(60 + (120 * y / height))
            b = int(85 + (80 * y / height))
            draw.line([(0, y), (width, y)], fill=(r, g, b))

    if add_shapes:
        random.seed(42)  # deterministic
        for _ in range(20):
            x1, y1 = random.randint(0, width), random.randint(0, height)
            x2, y2 = x1 + random.randint(20, 150), y1 + random.randint(20, 150)
            color = (random.randint(50, 255), random.randint(50, 255), random.randint(50, 255))
            shape_type = random.choice(["rect", "ellipse"])
            if shape_type == "rect":
                draw.rectangle([x1, y1, x2, y2], fill=color, outline=(255, 255, 255))
            else:
                draw.ellipse([x1, y1, x2, y2], fill=color, outline=(255, 255, 255))

    if add_noise:
        import numpy as np
        noise_arr = np.random.randint(0, 256, (height, width, 3), dtype=np.uint8)
        noise_img = Image.fromarray(noise_arr)
        img = Image.blend(img, noise_img, alpha=0.3)
        draw = ImageDraw.Draw(img)

    if text:
        draw.text((10, 10), text, fill=(255, 255, 255))

    buffer = io.BytesIO()
    save_kwargs = {"format": fmt}
    if fmt == "JPEG":
        save_kwargs["quality"] = 90
    img.save(buffer, **save_kwargs)
    return buffer.getvalue()


def generate_png_small() -> TestFile:
    """100×100 solid color with text overlay."""
    return TestFile(
        name="png_small.png",
        test_case_id="png_small",
        file_format="PNG",
        size_variant="small",
        content_type="image/png",
        file_bytes=_build_image(100, 100, "PNG", "Small Test"),
        description="100×100 PNG with solid background and text",
    )


def generate_png_medium() -> TestFile:
    """800×600 with geometric shapes."""
    return TestFile(
        name="png_medium.png",
        test_case_id="png_medium",
        file_format="PNG",
        size_variant="medium",
        content_type="image/png",
        file_bytes=_build_image(800, 600, "PNG", "Medium Test", add_shapes=True),
        description="800×600 PNG with shapes and text",
    )


def generate_png_large() -> TestFile:
    """4000×3000 high-res with gradient."""
    return TestFile(
        name="png_large.png",
        test_case_id="png_large",
        file_format="PNG",
        size_variant="large",
        content_type="image/png",
        file_bytes=_build_image(4000, 3000, "PNG", "Large Test", add_gradient=True, add_shapes=True),
        description="4000×3000 PNG with gradient and shapes",
    )


# ---------------------------------------------------------------------------
# JPG Generation
# ---------------------------------------------------------------------------

def generate_jpg_small() -> TestFile:
    """100×100 simple JPEG."""
    return TestFile(
        name="jpg_small.jpg",
        test_case_id="jpg_small",
        file_format="JPG",
        size_variant="small",
        content_type="image/jpeg",
        file_bytes=_build_image(100, 100, "JPEG", "Small JPEG"),
        description="100×100 JPEG with solid background",
    )


def generate_jpg_medium() -> TestFile:
    """1200×800 JPEG with shapes."""
    return TestFile(
        name="jpg_medium.jpg",
        test_case_id="jpg_medium",
        file_format="JPG",
        size_variant="medium",
        content_type="image/jpeg",
        file_bytes=_build_image(1200, 800, "JPEG", "Medium JPEG", add_shapes=True, add_gradient=True),
        description="1200×800 JPEG with gradient and shapes",
    )


def generate_jpg_large() -> TestFile:
    """4000×3000 JPEG with noise pattern."""
    try:
        data = _build_image(4000, 3000, "JPEG", "Large JPEG", add_gradient=True, add_noise=True)
    except ImportError:
        # numpy not available, fall back to gradient only
        data = _build_image(4000, 3000, "JPEG", "Large JPEG", add_gradient=True, add_shapes=True)
    return TestFile(
        name="jpg_large.jpg",
        test_case_id="jpg_large",
        file_format="JPG",
        size_variant="large",
        content_type="image/jpeg",
        file_bytes=data,
        description="4000×3000 JPEG with noise pattern",
    )


# ---------------------------------------------------------------------------
# PPTX Generation (requires python-pptx)
# ---------------------------------------------------------------------------

def _build_pptx(slides: list[tuple[str, str]]) -> bytes:
    """Build a PPTX from a list of (title, content) tuples."""
    if not HAS_PPTX:
        raise ImportError(
            "python-pptx is required for PPTX generation. "
            "Install with: pip install python-pptx"
        )

    prs = Presentation()

    for title_text, content_text in slides:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text

        # Set body text — handle case where placeholder might not exist
        try:
            body_placeholder = slide.placeholders[1]
            body_placeholder.text = content_text
        except (KeyError, IndexError):
            pass  # slide layout may not have a body placeholder

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def generate_pptx_small() -> TestFile:
    """1 slide, title only."""
    slides = _repeat_slides(1)
    return TestFile(
        name="pptx_small.pptx",
        test_case_id="pptx_small",
        file_format="PPTX",
        size_variant="small",
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        file_bytes=_build_pptx(slides),
        description="1-slide PPTX with title and content",
    )


def generate_pptx_medium() -> TestFile:
    """10 slides with bullet points."""
    slides = _repeat_slides(10)
    return TestFile(
        name="pptx_medium.pptx",
        test_case_id="pptx_medium",
        file_format="PPTX",
        size_variant="medium",
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        file_bytes=_build_pptx(slides),
        description="10-slide PPTX with varied content",
    )


def generate_pptx_large() -> TestFile:
    """50 slides with dense content."""
    slides = _repeat_slides(50)
    return TestFile(
        name="pptx_large.pptx",
        test_case_id="pptx_large",
        file_format="PPTX",
        size_variant="large",
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        file_bytes=_build_pptx(slides),
        description="50-slide PPTX with dense content",
    )


# ---------------------------------------------------------------------------
# Edge Case Files
# ---------------------------------------------------------------------------

def generate_edge_corrupt_pdf() -> TestFile:
    """Random bytes with PDF MIME type — should cause graceful extraction failure."""
    random.seed(99)
    corrupt_bytes = bytes(random.getrandbits(8) for _ in range(2048))
    return TestFile(
        name="edge_corrupt.pdf",
        test_case_id="edge_corrupt_pdf",
        file_format="PDF",
        size_variant="edge",
        content_type="application/pdf",
        file_bytes=corrupt_bytes,
        description="2 KB of random bytes with PDF MIME type — tests graceful failure",
        is_edge_case=True,
        expected_failure=True,
    )


def generate_edge_empty_docx() -> TestFile:
    """0-byte file with DOCX content type — should fail gracefully."""
    return TestFile(
        name="edge_empty.docx",
        test_case_id="edge_empty_docx",
        file_format="DOCX",
        size_variant="edge",
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        file_bytes=b"",
        description="0-byte file with DOCX MIME type — tests empty file handling",
        is_edge_case=True,
        expected_failure=True,
    )


def generate_edge_large_single_page_pdf() -> TestFile:
    """Single very tall page with massive text — tests chunking under extreme conditions."""
    if not HAS_REPORTLAB:
        raise ImportError("reportlab required for edge_large_single_page_pdf")

    # Create a very tall single page with ~200 KB of text
    page_width = 8.5 * inch
    page_height = 500 * inch  # extremely tall page

    paragraphs = _repeat_paragraphs(800)

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=(page_width, page_height))
    styles = getSampleStyleSheet()
    normal = styles["Normal"]

    story = []
    for p in paragraphs:
        story.append(Paragraph(p, normal))
        story.append(Spacer(1, 6))

    doc.build(story)
    data = buffer.getvalue()

    return TestFile(
        name="edge_large_single_page.pdf",
        test_case_id="edge_large_single_page_pdf",
        file_format="PDF",
        size_variant="edge",
        content_type="application/pdf",
        file_bytes=data,
        description=(
            f"Single very-tall-page PDF ({len(data):,} bytes) with ~800 paragraphs "
            "— tests chunking behavior on a single source page"
        ),
        is_edge_case=True,
        expected_failure=False,
    )


def generate_edge_no_text_png() -> TestFile:
    """Gradient-only image with no readable text — image pipeline should still succeed."""
    img_bytes = _build_image(
        800, 600, "PNG", text="", add_gradient=True, add_shapes=False
    )
    return TestFile(
        name="edge_no_text.png",
        test_case_id="edge_no_text_png",
        file_format="PNG",
        size_variant="edge",
        content_type="image/png",
        file_bytes=img_bytes,
        description="800×600 gradient PNG with no text — tests image-only pipeline path",
        is_edge_case=True,
        expected_failure=False,
    )


def generate_edge_duplicate_txt() -> list[TestFile]:
    """Two identical text files — tests that processing the same file twice works."""
    text = "This is a duplicate upload test. " * 20
    file_bytes = text.encode("utf-8")
    base = TestFile(
        name="edge_duplicate.txt",
        test_case_id="edge_duplicate_upload",
        file_format="TXT",
        size_variant="edge",
        content_type="text/plain",
        file_bytes=file_bytes,
        description="Duplicate upload test — same file processed twice for idempotency check",
        is_edge_case=True,
        expected_failure=False,
    )
    # Return two copies — the test runner handles them as duplicate uploads
    return [base]


# ---------------------------------------------------------------------------
# Main generator
# ---------------------------------------------------------------------------

# Ordered list of standard generators
_STANDARD_GENERATORS = [
    generate_pdf_small,
    generate_pdf_medium,
    generate_pdf_large,
    generate_docx_small,
    generate_docx_medium,
    generate_docx_large,
    generate_txt_small,
    generate_txt_medium,
    generate_txt_large,
    generate_png_small,
    generate_png_medium,
    generate_png_large,
    generate_jpg_small,
    generate_jpg_medium,
    generate_jpg_large,
    generate_pptx_small,
    generate_pptx_medium,
    generate_pptx_large,
]

_EDGE_GENERATORS = [
    generate_edge_corrupt_pdf,
    generate_edge_empty_docx,
    generate_edge_large_single_page_pdf,
    generate_edge_no_text_png,
    generate_edge_duplicate_txt,
]


def generate_all() -> list[TestFile]:
    """
    Generate all test files: 18 standard + 5 edge cases = 23 test cases.

    Skips formats whose dependencies are not installed (with a warning),
    rather than crashing the entire test run.
    """
    files: list[TestFile] = []

    # --- Standard test matrix ---
    for gen in _STANDARD_GENERATORS:
        try:
            files.append(gen())
        except ImportError as e:
            logger.warning("Skipping %s: %s", gen.__name__, e)
        except Exception as e:
            logger.error("Failed to generate %s: %s", gen.__name__, e)

    # --- Edge cases ---
    for gen in _EDGE_GENERATORS:
        try:
            result = gen()
            if isinstance(result, list):
                files.extend(result)
            else:
                files.append(result)
        except ImportError as e:
            logger.warning("Skipping edge case %s: %s", gen.__name__, e)
        except Exception as e:
            logger.error("Failed to generate edge case %s: %s", gen.__name__, e)

    return files


def check_dependencies() -> dict[str, bool]:
    """Check which optional dependencies are available."""
    return {
        "reportlab (PDF)": HAS_REPORTLAB,
        "python-docx (DOCX)": HAS_DOCX,
        "python-pptx (PPTX)": HAS_PPTX,
        "Pillow (PNG/JPG)": True,  # always available via requirements.txt
    }


# ---------------------------------------------------------------------------
# CLI preview
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    print("=" * 65)
    print("Test File Generator — Dependency Check")
    print("=" * 65)
    for dep, available in check_dependencies().items():
        icon = "✅" if available else "❌"
        print(f"  {icon}  {dep}")

    print(f"\n{'=' * 65}")
    print("Generating test files...")
    print("=" * 65)

    test_files = generate_all()

    print(f"\n{'#':<4} {'test_case_id':<30} {'format':<6} {'size':<8} {'bytes':>12}  {'description'}")
    print("-" * 100)
    for i, f in enumerate(test_files, 1):
        edge = " [EDGE]" if f.is_edge_case else ""
        fail = " [EXPECT FAIL]" if f.expected_failure else ""
        print(f"{i:<4} {f.test_case_id:<30} {f.file_format:<6} {f.size_variant:<8} "
              f"{f.file_size_bytes:>12,}  {f.description[:50]}{edge}{fail}")

    total_bytes = sum(f.file_size_bytes for f in test_files)
    print(f"\nTotal: {len(test_files)} files, {total_bytes:,} bytes ({total_bytes / 1024 / 1024:.1f} MB)")
